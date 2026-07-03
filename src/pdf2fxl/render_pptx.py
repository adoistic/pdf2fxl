from __future__ import annotations
from pathlib import Path
from typing import List, Tuple
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .models import Page

_EMU_PER_IN = 914400
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


def _hex_to_rgb(c: str) -> RGBColor:
    c = c.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def write_pptx(pages: List[Page], out_path: Path, aspect: Tuple[int, int] = (4, 3),
               font_name: str = "Noto Serif", slide_height_in: float = 7.5) -> Path:
    prs = Presentation()
    prs.slide_height = Emu(int(slide_height_in * _EMU_PER_IN))
    if pages:
        pw0, ph0 = pages[0].page_size_px
        page_ratio = pw0 / ph0
    else:
        aw, ah = aspect
        page_ratio = aw / ah
    prs.slide_width = Emu(int(slide_height_in * page_ratio * _EMU_PER_IN))
    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)
    slide_h_pt = slide_height_in * 72
    blank = prs.slide_layouts[6]

    for page in pages:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(page.background, 0, 0, width=slide_w, height=slide_h)
        pw, ph = page.page_size_px
        for b in sorted(page.blocks, key=lambda z: z.reading_order):
            x, y, w, h = b.bbox
            box = slide.shapes.add_textbox(
                Emu(int(x / pw * slide_w)), Emu(int(y / ph * slide_h)),
                Emu(int(w / pw * slide_w)), Emu(int(h / ph * slide_h)))
            tf = box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            para = tf.paragraphs[0]
            para.alignment = _ALIGN.get(b.align, PP_ALIGN.LEFT)
            run = para.add_run()
            run.text = b.text
            run.font.name = font_name
            run.font.size = Pt(round(b.font_px / ph * slide_h_pt, 1))
            run.font.color.rgb = _hex_to_rgb(b.color)
    prs.save(str(out_path))
    return Path(out_path)
