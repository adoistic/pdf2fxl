from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from pdf2fxl.models import Page, Block
from pdf2fxl.render_pptx import write_pptx

def _one_page(tmp_path):
    bg = tmp_path / "page-00-clean.png"
    Image.new("RGB", (400, 300), (255, 255, 255)).save(bg)
    return Page(index=0, page_size_px=(400, 300), background=str(bg), original=str(bg),
                blocks=[Block(type="text", bbox=(40, 60, 200, 90), text="Hello deck",
                              font_px=24, color="#222222", align="left", reading_order=0)])

def test_pptx_slide_picture_and_textbox(tmp_path):
    out = tmp_path / "book.pptx"
    write_pptx([_one_page(tmp_path)], out, aspect=(4, 3))
    prs = Presentation(str(out))
    assert prs.slide_width == Emu(9144000) and prs.slide_height == Emu(6858000)
    slide = prs.slides[0]
    pics = [s for s in slide.shapes if s.shape_type == 13]        # PICTURE
    boxes = [s for s in slide.shapes if s.has_text_frame]
    assert len(pics) == 1 and len(boxes) == 1
    assert boxes[0].text_frame.text == "Hello deck"
    # left = 40/400 * 9144000 = 914400
    assert abs(boxes[0].left - 914400) < 2000

def test_pptx_slide_matches_page_aspect(tmp_path):
    # a 3:2 page (600x400) must yield a 3:2 slide, not 4:3
    bg = tmp_path / "p.png"; Image.new("RGB", (600, 400), (255,255,255)).save(bg)
    page = Page(index=0, page_size_px=(600, 400), background=str(bg), original=str(bg),
                blocks=[Block(type="text", bbox=(60, 40, 200, 60), text="x",
                              font_px=20, color="#000000", align="left", reading_order=0)])
    out = tmp_path / "d.pptx"
    write_pptx([page], out)
    prs = Presentation(str(out))
    ratio = prs.slide_width / prs.slide_height
    assert abs(ratio - 600/400) < 0.001    # 1.5, i.e. 3:2 — not 1.333
